#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
EDR Capa 專案進度回報簡報生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

def create_progress_presentation():
    """建立 EDR Capa 進度回報簡報"""

    # 建立簡報物件
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 設定預設字型
    def set_text_frame_format(text_frame, font_size=14, bold=False):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = '微軟正黑體'
                run.font.size = Pt(font_size)
                run.font.bold = bold

    # ========== 投影片 1: 封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版型

    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "EDR Capa 支援專案"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "雙週進度回報"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(64, 64, 64)

    # 日期
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"回報日期：{datetime.now().strftime('%Y年%m月%d日')}"
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(128, 128, 128)

    # ========== 投影片 2: 專案概述 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 標題和內容
    title = slide.shapes.title
    title.text = "專案概述"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "目標"

    p = tf.add_paragraph()
    p.text = "將 Mandiant Capa 規則引擎整合至 EDR 系統"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "實現即時威脅偵測與行為分析"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "即時威脅阻擋與防禦"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "自動化響應與處置"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "核心模組"
    p.level = 0

    modules = [
        "模組 1：規則解析器（12 任務）",
        "模組 2：事件抽象層（12 任務）",
        "模組 3：匹配引擎（31 任務）",
        "模組 4：即時決策與阻擋引擎（8 任務）",
        "模組 5：響應動作執行器（15 任務）",
        "模組 6：告警生成與管理（9 任務）",
        "模組 7：效能監控（9 任務）",
        "模組 8：整合與配置（10 任務）",
        "模組 9：引擎 API（7 任務）"
    ]

    for module in modules:
        p = tf.add_paragraph()
        p.text = module
        p.level = 1

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = '微軟正黑體'
            run.font.size = Pt(16)

    # ========== 投影片 3: 整體進度概覽 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 標題
    title = slide.shapes.title
    title.text = "整體進度概覽"
    title.text_frame.paragraphs[0].font.size = Pt(32)

    # 統計表格
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(3.5)

    shape = slide.shapes.add_table(6, 4, left, top, width, height)
    table = shape.table

    # 設定表頭
    headers = ['項目', '總數', '已完成', '完成率']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 填入數據
    data = [
        ['P0（關鍵）任務', '55', '0', '0%'],
        ['P1（重要）任務', '45', '0', '0%'],
        ['P2（一般）任務', '13', '0', '0%'],
        ['總計任務數', '113', '0', '0%'],
        ['預估完成時間', '7-8個月', '-', '-']
    ]

    for i, row in enumerate(data, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if j == 0:
                cell.text_frame.paragraphs[0].font.bold = True

    # ========== 投影片 4-10: 各模組詳細進度 ==========
    modules_detail = [
        {
            'name': '模組 1：規則解析器',
            'total': 12,
            'p0': 9,
            'p1': 2,
            'p2': 1,
            'tasks': [
                'YAML 函式庫整合',
                'Meta 資料結構定義',
                'FeatureNode 基礎介面',
                '邏輯節點類別（AND/OR/NOT/COUNT）',
                '基本陳述節點（API/String/Bytes/Number）',
                '進階陳述節點（OS/Arch/Import/Export）',
                'Meta 區塊解析器',
                'Features 區塊遞迴解析器',
                '規則驗證器',
                'RuleParser 主類別',
                '規則載入管理',
                '規則快取機制'
            ]
        },
        {
            'name': '模組 2：事件抽象層',
            'total': 12,
            'p0': 7,
            'p1': 4,
            'p2': 1,
            'tasks': [
                'ApiCallEvent 結構定義',
                'ProcessMemoryRegion 結構定義',
                'SystemInfo 結構定義',
                'ProcessInfo 結構定義',
                '其他事件結構定義',
                'API 呼叫事件正規化器',
                '程序事件正規化器',
                '記憶體事件正規化器',
                '檔案系統事件正規化器',
                '登錄檔事件正規化器',
                '事件佇列實作',
                '事件快取機制'
            ]
        },
        {
            'name': '模組 3：匹配引擎（1/2）',
            'total': 31,
            'p0': 15,
            'p1': 13,
            'p2': 3,
            'tasks': [
                'EvaluationContext 類別',
                'getApiCallHistory 方法',
                'getProcessMemory 方法',
                'getSystemInfo 方法',
                'getRuleResult 方法',
                'AndNode 評估器',
                'OrNode 評估器',
                'NotNode 評估器',
                'CountNode 評估器',
                'ApiStatement 評估器',
                'StringStatement 評估器',
                'SubstringStatement 評估器'
            ]
        },
        {
            'name': '模組 3：匹配引擎（2/2）',
            'total': 31,
            'p0': 15,
            'p1': 13,
            'p2': 3,
            'tasks': [
                'BytesStatement 評估器',
                'NumberStatement 評估器',
                'OsStatement 評估器',
                'ArchStatement 評估器',
                'FormatStatement 評估器',
                'CharacteristicStatement 評估器',
                'ImportStatement 評估器',
                'ExportStatement 評估器',
                'SectionStatement 評估器',
                'MatchStatement 評估器',
                'PE 檔案解析器（標頭/區段/匯入/匯出）',
                '評估快取、記憶體快取、事件預過濾器'
            ]
        },
        {
            'name': '模組 4：即時決策與阻擋引擎',
            'total': 8,
            'p0': 4,
            'p1': 3,
            'p2': 1,
            'tasks': [
                '威脅等級評估器',
                '行為上下文分析器',
                '阻擋策略決策引擎',
                '誤報風險評估',
                '白名單管理',
                '阻擋配置管理',
                '阻擋模式切換',
                '決策日誌記錄'
            ]
        },
        {
            'name': '模組 5：響應動作執行器（1/2）',
            'total': 15,
            'p0': 9,
            'p1': 4,
            'p2': 2,
            'tasks': [
                'Kernel Hook 機制',
                'User-mode Hook 機制',
                'API 阻擋邏輯',
                '關鍵 API 攔截',
                '程序終止功能',
                '程序掛起功能',
                '檔案隔離功能',
                '檔案還原功能'
            ]
        },
        {
            'name': '模組 5：響應動作執行器（2/2）',
            'total': 15,
            'p0': 9,
            'p1': 4,
            'p2': 2,
            'tasks': [
                '隔離區清理',
                '惡意記憶體清除',
                '記憶體保護',
                '登錄檔還原',
                '檔案系統還原',
                '網路連線阻斷',
                '程序網路隔離'
            ]
        },
        {
            'name': '模組 6：告警生成與管理',
            'total': 9,
            'p0': 2,
            'p1': 6,
            'p2': 1,
            'tasks': [
                'CapaMatchAlert 類別',
                '告警生成器',
                '觸發路徑追蹤',
                '告警去重機制',
                '告警聚合',
                'HTTP 告警傳送器',
                '告警重試機制',
                '告警本地快取',
                '告警批次傳送'
            ]
        },
        {
            'name': '模組 7：效能監控',
            'total': 9,
            'p0': 0,
            'p1': 7,
            'p2': 2,
            'tasks': [
                'CPU 使用率監控',
                '記憶體使用監控',
                '延遲監控（P50/P95/P99）',
                '吞吐量監控',
                '快取命中率監控',
                '效能日誌輸出',
                'Prometheus 匯出器',
                '健康檢查端點',
                '效能異常告警'
            ]
        },
        {
            'name': '模組 8：整合與配置',
            'total': 10,
            'p0': 5,
            'p1': 3,
            'p2': 2,
            'tasks': [
                'EDR 事件接收器',
                '主工作流程實作',
                '規則評估排程器',
                '配置檔案解析器',
                'EngineConfig 類別',
                '配置熱更新',
                '日誌記錄器',
                '日誌輪轉',
                '日誌過濾'
            ]
        },
        {
            'name': '模組 9：匹配引擎 API',
            'total': 7,
            'p0': 4,
            'p1': 3,
            'p2': 0,
            'tasks': [
                'MatchingEngine 類別',
                'loadRules 方法',
                'processEvent 方法',
                'processBatch 方法',
                'getStats 方法',
                'updateRules 方法',
                'shutdown 方法'
            ]
        }
    ]

    for module in modules_detail:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = module['name']
        title.text_frame.paragraphs[0].font.size = Pt(28)

        # 模組資訊
        info_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
        info_frame = info_box.text_frame
        info_frame.text = f"總任務數：{module['total']} | P0: {module['p0']} | P1: {module['p1']} | P2: {module['p2']}"
        info_para = info_frame.paragraphs[0]
        info_para.font.size = Pt(16)
        info_para.font.bold = True
        info_para.font.color.rgb = RGBColor(0, 51, 102)

        # 任務清單
        task_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4.5))
        task_frame = task_box.text_frame
        task_frame.word_wrap = True

        for i, task in enumerate(module['tasks'], 1):
            p = task_frame.add_paragraph() if i > 1 else task_frame.paragraphs[0]
            p.text = f"□ {task}"
            p.font.size = Pt(14)
            p.font.name = '微軟正黑體'
            p.space_after = Pt(8)

    # ========== 投影片 11: 本期（雙週）計畫 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "本期計畫（雙週）"
    title.text_frame.paragraphs[0].font.size = Pt(32)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "預計完成任務"

    planned_tasks = [
        "□ PARSER-001: 選擇並整合 YAML 函式庫",
        "□ PARSER-002: 定義 Meta 資料結構",
        "□ PARSER-003: 定義 FeatureNode 基礎介面",
        "□ EVENT-001: 定義 ApiCallEvent 結構",
        "□ EVENT-002: 定義 ProcessMemoryRegion 結構",
        "□ 建立專案基礎架構和建置系統"
    ]

    for task in planned_tasks:
        p = tf.add_paragraph()
        p.text = task
        p.level = 1
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "風險與挑戰"
    p.level = 0

    risks = [
        "YAML 函式庫選擇需要評估效能和安全性",
        "資料結構設計需要考慮未來擴展性"
    ]

    for risk in risks:
        p = tf.add_paragraph()
        p.text = risk
        p.level = 1
        p.font.size = Pt(16)

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = '微軟正黑體'

    # ========== 投影片 12: 下期預覽 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "下期預覽"
    title.text_frame.paragraphs[0].font.size = Pt(32)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "計畫目標"

    next_tasks = [
        "完成規則解析器的核心實作",
        "實作邏輯節點類別（AND/OR/NOT/COUNT）",
        "實作基本陳述節點（API/String/Bytes/Number）",
        "開始事件抽象層的實作",
        "建立初步的單元測試框架"
    ]

    for task in next_tasks:
        p = tf.add_paragraph()
        p.text = task
        p.level = 1
        p.font.size = Pt(18)

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = '微軟正黑體'

    # ========== 投影片 13: 問題與討論 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "問題與討論"
    title.text_frame.paragraphs[0].font.size = Pt(32)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "需要討論的議題"

    discussions = [
        "YAML 解析函式庫選型建議",
        "資料結構設計審查",
        "開發環境與工具鏈確認",
        "測試策略討論",
        "資源需求確認"
    ]

    for item in discussions:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "需要的支援"
    p.level = 0

    supports = [
        "EDR 系統 API 文件",
        "測試環境存取權限",
        "Capa 規則範例"
    ]

    for item in supports:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(20)

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = '微軟正黑體'

    # ========== 投影片 14: Q&A ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    qa_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Q & A"
    qa_para = qa_frame.paragraphs[0]
    qa_para.alignment = PP_ALIGN.CENTER
    qa_para.font.size = Pt(60)
    qa_para.font.bold = True
    qa_para.font.color.rgb = RGBColor(0, 51, 102)

    # 儲存簡報
    filename = f'EDR-Capa-進度回報-{datetime.now().strftime("%Y%m%d")}.pptx'
    prs.save(filename)
    print(f"[OK] 簡報已成功建立：{filename}")
    return filename

if __name__ == '__main__':
    create_progress_presentation()