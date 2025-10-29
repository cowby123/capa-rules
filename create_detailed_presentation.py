#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
EDR Capa 專案詳細規格簡報生成器
包含完整的專案規格、架構設計、開發計畫等資訊
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

def add_title_slide(prs, title_text, subtitle_text):
    """添加標題投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_text
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(64, 64, 64)

    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"{datetime.now().strftime('%Y年%m月%d日')}"
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(128, 128, 128)

def add_content_slide(prs, title_text, content_dict):
    """添加內容投影片（標題 + 多段內容）"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    first_section = True
    for section_title, items in content_dict.items():
        if not first_section:
            p = tf.add_paragraph()
            p.text = ""

        p = tf.add_paragraph() if not first_section else tf.paragraphs[0]
        p.text = section_title
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = '微軟正黑體'

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(14)
            p.font.name = '微軟正黑體'

        first_section = False

def add_table_slide(prs, title_text, headers, data):
    """添加表格投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)

    rows = len(data) + 1
    cols = len(headers)

    left = Inches(0.5) if cols > 5 else Inches(1.5)
    top = Inches(2)
    width = Inches(9) if cols > 5 else Inches(7)
    height = Inches(4)

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    # 設定表頭
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 填入數據
    for i, row in enumerate(data, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.name = '微軟正黑體'
            if j == 0:
                cell.text_frame.paragraphs[0].font.bold = True

def create_detailed_presentation():
    """建立詳細的 EDR Capa 專案規格簡報"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== 投影片 1: 封面 ==========
    add_title_slide(prs, "EDR Capa 支援專案", "詳細規格與開發計畫")

    # ========== 投影片 2: 目錄 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "目錄"
    title.text_frame.paragraphs[0].font.size = Pt(36)

    content = slide.placeholders[1]
    tf = content.text_frame

    toc_items = [
        "1. 專案概述與目標",
        "2. 系統架構設計",
        "3. 核心功能規格",
        "4. 即時阻擋機制",
        "5. 技術規格詳述",
        "6. 開發計畫與時程",
        "7. 模組詳細設計（9個模組）",
        "8. 效能需求與指標",
        "9. 風險評估與對策",
        "10. 資源需求",
        "11. 進度追蹤"
    ]

    tf.text = toc_items[0]
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.name = '微軟正黑體'

    for item in toc_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.name = '微軟正黑體'
        p.space_after = Pt(10)

    # ========== 投影片 3: 專案概述 ==========
    add_content_slide(prs, "專案概述", {
        "專案背景": [
            "Mandiant Capa：開源的惡意程式能力檢測工具",
            "5000+ YAML 規則，涵蓋 ATT&CK 和 MBC 框架",
            "目前僅支援靜態分析，需要轉換為動態監控"
        ],
        "專案目標": [
            "將 Capa 規則引擎整合至 EDR 系統",
            "實現即時威脅偵測與行為分析",
            "即時威脅阻擋與自動化響應",
            "提供完整的威脅情報映射（ATT&CK/MBC）"
        ],
        "預期成果": [
            "支援 5000+ 規則的動態執行",
            "事件處理延遲 < 50ms (P95)",
            "每秒處理 ≥ 15,000 事件（4 核心）",
            "CPU 使用率 < 30%（高負載）"
        ]
    })

    # ========== 投影片 4: 系統架構 - 整體視圖 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "系統架構 - 整體視圖"
    title.text_frame.paragraphs[0].font.size = Pt(32)

    arch_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    arch_frame = arch_box.text_frame
    arch_frame.word_wrap = True

    architecture_text = """┌─────────────────────────────────────────────────────────┐
│                     EDR 事件源                          │
│  (程序監控 / API Hook / 檔案監控 / 網路監控 / ETW)    │
└──────────────────────┬──────────────────────────────────┘
                       ↓
┌─────────────────────────────────────────────────────────┐
│              事件抽象層（Event Abstraction）            │
│        事件標準化 / 事件佇列 / 事件快取                │
└──────────────────────┬──────────────────────────────────┘
                       ↓
┌─────────────────────────────────────────────────────────┐
│         規則匹配引擎（Capa Rule Matching Engine）       │
│    規則解析 / 特徵評估 / 邏輯運算 / PE 解析            │
└──────────────────────┬──────────────────────────────────┘
                       ↓
┌─────────────────────────────────────────────────────────┐
│      即時決策引擎（Real-time Decision Engine）         │
│  威脅評分 / 策略選擇 / 白名單檢查 / 誤報評估           │
└──────────┬────────────────────────────┬─────────────────┘
           ↓                            ↓
┌─────────────────────┐    ┌──────────────────────────────┐
│  告警生成與管理      │    │   響應動作執行器              │
│  Alert Generation   │    │   Response Executor          │
│  • 告警生成         │    │   • API 攔截與阻擋           │
│  • 告警去重         │    │   • 程序終止/掛起            │
│  • 告警傳送         │    │   • 檔案隔離與還原           │
└─────────────────────┘    │   • 記憶體清除               │
                           │   • 系統還原                 │
                           │   • 網路隔離                 │
                           └──────────────────────────────┘"""

    arch_frame.text = architecture_text
    for paragraph in arch_frame.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.name = 'Consolas'

    # ========== 投影片 5: 系統架構 - 資料流 ==========
    add_content_slide(prs, "系統架構 - 資料流程", {
        "1. 事件收集": [
            "從 EDR 核心驅動接收原始事件",
            "支援 API 呼叫、程序操作、檔案操作、網路活動等",
            "事件頻率：高峰可達 50,000 events/sec"
        ],
        "2. 事件標準化": [
            "將 EDR 事件轉換為 Capa 可理解的格式",
            "提取關鍵特徵（API名稱、參數、返回值等）",
            "建立程序上下文（PID、路徑、記憶體映像）"
        ],
        "3. 規則評估": [
            "載入並快取 Capa 規則（5000+ 規則）",
            "並行評估匹配條件（AND/OR/NOT/COUNT）",
            "短路評估以提升效能"
        ],
        "4. 威脅決策": [
            "計算威脅分數（0-100）",
            "選擇響應策略（Kill/Block/Suspend/Monitor/Quarantine）",
            "考慮白名單和組織例外"
        ],
        "5. 執行響應": [
            "根據策略執行阻擋動作",
            "記錄完整的操作日誌",
            "生成告警並傳送至 SIEM"
        ]
    })

    # ========== 投影片 6: 核心功能 - 即時阻擋機制 ==========
    add_content_slide(prs, "核心功能 - 即時阻擋機制", {
        "威脅等級定義": [
            "Critical (90-100)：勒索軟體、記憶體注入、批量刪除",
            "High (70-89)：停用防毒、修改關鍵登錄檔、建立服務",
            "Medium (40-69)：異常但可能合法的行為",
            "Low (0-39)：低風險探測行為"
        ],
        "五種阻擋策略": [
            "Kill：立即終止程序（Critical 威脅）",
            "Block：阻止特定 API 但不終止程序（High 威脅）",
            "Suspend：暫停程序等待分析（高誤報風險）",
            "Monitor：允許執行但加強監控（中低威脅）",
            "Quarantine：檔案隔離（高危檔案操作）"
        ],
        "決策因素": [
            "規則元資料（ATT&CK、MBC、severity）",
            "行為上下文（程序路徑、簽章、父程序）",
            "行為序列和時間關聯性",
            "白名單和組織例外規則"
        ]
    })

    # ========== 投影片 7: 響應動作執行器詳述 ==========
    add_content_slide(prs, "響應動作執行器 - 技術細節", {
        "API 攔截技術": [
            "Kernel Hook：SSDT Hook、IRP Hook（核心層攔截）",
            "User-mode Hook：IAT Hook、Inline Hook（使用者層攔截）",
            "支援同步阻擋、非同步阻擋、欺騙模式"
        ],
        "程序控制": [
            "強制終止：TerminateProcess、NtTerminateProcess、ZwTerminateProcess",
            "程序掛起：SuspendThread、NtSuspendProcess",
            "子程序處理：遞迴終止所有子程序"
        ],
        "檔案隔離": [
            "移動檔案至隔離區（隱藏分區）",
            "AES-256 加密存儲",
            "保留元資料：原始路徑、時間戳、觸發規則、SHA256",
            "支援授權還原"
        ],
        "系統還原": [
            "登錄檔快照與還原（RegSaveKey、RegRestoreKey）",
            "檔案系統還原（VSS 影子複製）",
            "記錄所有變更以支援審計"
        ]
    })

    # ========== 投影片 8: 技術規格 - 資料結構 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "技術規格 - 核心資料結構"
    title.text_frame.paragraphs[0].font.size = Pt(32)

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True

    code_text = """// 規則物件
class Rule {
    Meta meta;              // 規則元資料
    FeatureNode* features;  // 特徵樹根節點
    bool evaluate(EvaluationContext& ctx) const;
};

// 評估上下文
class EvaluationContext {
    ProcessInfo processInfo;                    // 程序資訊
    std::vector<ApiCallEvent> apiCalls;        // API 呼叫歷史
    std::map<std::string, std::vector<std::string>> imports;  // 匯入函式
    std::map<FeatureNode*, bool> cache;        // 評估快取
};

// 告警物件
struct CapaMatchAlert {
    std::string ruleName;                      // 規則名稱
    ProcessInfo processInfo;                   // 程序資訊
    std::string threatLevel;                   // 威脅等級
    std::string recommendedAction;             // 建議動作
    std::vector<std::string> attackTechniques; // ATT&CK 技術
    std::string timestamp;                     // 時間戳
};

// 決策結果
struct DecisionResult {
    int threatScore;           // 威脅分數 0-100
    BlockingStrategy strategy; // 阻擋策略
    std::string reason;        // 決策理由
    bool whitelisted;          // 是否白名單
};"""

    code_frame.text = code_text
    for paragraph in code_frame.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.name = 'Consolas'

    # ========== 投影片 9: 開發計畫 - 階段劃分 ==========
    add_content_slide(prs, "開發計畫 - 三階段實施", {
        "第一階段：核心基礎（P0，4-5個月）": [
            "規則解析器：YAML 解析、規則驗證、規則載入",
            "事件抽象層：事件標準化、佇列管理",
            "匹配引擎核心：邏輯節點、基本特徵評估",
            "決策與阻擋引擎：威脅評估、策略選擇、決策管理",
            "響應執行器：API 攔截、程序終止、檔案隔離",
            "整合與 API：EDR 整合、配置管理"
        ],
        "第二階段：進階功能（P1，3個月）": [
            "進階特徵評估：PE 解析、Section、Match 等",
            "快取與最佳化：多層快取、記憶體快取",
            "決策引擎進階：誤報評估、白名單管理",
            "響應執行器進階：記憶體清除、系統還原、網路隔離",
            "效能監控：CPU、記憶體、延遲監控",
            "告警管理：去重、批次傳送、重試"
        ],
        "第三階段：優化與增強（P2，1個月）": [
            "規則快取、事件預過濾器",
            "Prometheus 整合、效能異常告警",
            "配置熱更新",
            "隔離區管理、記憶體保護"
        ]
    })

    # ========== 投影片 10: 時程規劃表 ==========
    headers = ['階段', '時長', '主要交付物', '關鍵里程碑']
    data = [
        ['階段一', '4-5個月', '核心引擎、基礎阻擋', '完成 P0 任務（55個）'],
        ['階段二', '3個月', '進階功能、效能優化', '完成 P1 任務（45個）'],
        ['階段三', '1個月', '最終優化、文件', '完成 P2 任務（13個）'],
        ['總計', '8-9個月', '完整系統', '113 個任務']
    ]
    add_table_slide(prs, "時程規劃", headers, data)

    # ========== 投影片 11-19: 九個模組詳細設計 ==========
    modules_detail = [
        {
            'name': '模組 1：規則解析器',
            'total': 12, 'p0': 9, 'p1': 2, 'p2': 1,
            'description': '負責解析 Capa YAML 規則，建立規則物件樹',
            'key_components': [
                'YAML 解析引擎（yaml-cpp 或 RapidYAML）',
                '規則資料結構（Meta、FeatureNode）',
                'Meta 區塊解析器（名稱、作者、ATT&CK、MBC）',
                'Features 區塊遞迴解析器（AND/OR/NOT/COUNT）',
                '規則驗證器（語法檢查、依賴檢查）',
                '規則快取機制（記憶體快取）'
            ],
            'challenges': [
                'YAML 解析效能（5000+ 規則）',
                '規則依賴關係處理',
                '錯誤處理與診斷'
            ]
        },
        {
            'name': '模組 2：事件抽象層',
            'total': 12, 'p0': 7, 'p1': 4, 'p2': 1,
            'description': '將 EDR 原始事件標準化為 Capa 可處理的格式',
            'key_components': [
                '事件標準化器（EDR → Capa 格式轉換）',
                '事件佇列管理（多生產者/多消費者）',
                '事件快取（近期事件查詢）',
                '程序上下文管理（記憶體映像、匯入表）',
                '支援多種事件類型（API、檔案、網路、登錄檔）'
            ],
            'challenges': [
                '高頻事件處理（峰值 50k events/sec）',
                '記憶體使用控制',
                '事件順序保證'
            ]
        },
        {
            'name': '模組 3：匹配引擎',
            'total': 31, 'p0': 15, 'p1': 13, 'p2': 3,
            'description': '執行規則評估，匹配惡意行為模式',
            'key_components': [
                '評估上下文管理（程序資訊、API 歷史）',
                '邏輯節點評估（AND/OR/NOT/COUNT、短路）',
                '基本特徵評估（API/String/Bytes/Number）',
                '進階特徵評估（OS/Arch/Import/Export/Section）',
                'PE 檔案解析器（標頭/區段/匯入/匯出表）',
                '多層快取（評估快取、記憶體快取）'
            ],
            'challenges': [
                '評估延遲要求（<5ms per rule）',
                'PE 解析效能',
                '快取一致性'
            ]
        },
        {
            'name': '模組 4：即時決策與阻擋引擎',
            'total': 8, 'p0': 4, 'p1': 3, 'p2': 1,
            'description': '評估威脅等級並選擇適當的響應策略',
            'key_components': [
                '威脅等級評估器（計算威脅分數 0-100）',
                '行為上下文分析器（程序、簽章、序列）',
                '阻擋策略決策引擎（5種策略選擇）',
                '誤報風險評估',
                '白名單管理（路徑、簽章、規則例外）',
                '阻擋配置管理（全域開關、模式切換）'
            ],
            'challenges': [
                '誤報率控制（False Positive < 1%）',
                '決策延遲最小化',
                '白名單維護'
            ]
        },
        {
            'name': '模組 5：響應動作執行器',
            'total': 15, 'p0': 9, 'p1': 4, 'p2': 2,
            'description': '執行實際的阻擋和響應動作',
            'key_components': [
                'API 攔截（Kernel Hook、User-mode Hook）',
                '程序終止/掛起（含子程序處理）',
                '檔案隔離（AES-256 加密、元資料保存）',
                '檔案還原（解密、授權檢查）',
                '記憶體清除（Shellcode 清除、記憶體保護）',
                '系統還原（登錄檔、檔案系統）',
                '網路隔離（連線阻斷、防火牆規則）'
            ],
            'challenges': [
                'Hook 穩定性和安全性',
                '系統相容性（不同 Windows 版本）',
                '還原操作的可靠性'
            ]
        },
        {
            'name': '模組 6：告警生成與管理',
            'total': 9, 'p0': 2, 'p1': 6, 'p2': 1,
            'description': '生成、去重、傳送威脅告警',
            'key_components': [
                'CapaMatchAlert 類別（完整告警資訊）',
                '告警生成器（自動填充元資料）',
                '觸發路徑追蹤（解釋告警原因）',
                '告警去重機制（時間視窗、雜湊表）',
                'HTTP/HTTPS 傳送（TLS 1.3）',
                '重試與本地快取'
            ],
            'challenges': [
                '告警風暴預防',
                '傳送可靠性',
                '本地快取管理'
            ]
        },
        {
            'name': '模組 7：效能監控',
            'total': 9, 'p0': 0, 'p1': 7, 'p2': 2,
            'description': '收集並報告引擎效能指標',
            'key_components': [
                'CPU/記憶體使用率監控',
                '延遲監控（P50/P95/P99）',
                '吞吐量和快取命中率',
                '效能日誌輸出',
                'Prometheus 匯出器',
                '健康檢查端點',
                '效能異常告警'
            ],
            'challenges': [
                '監控開銷控制',
                '指標準確性',
                '異常檢測閾值設定'
            ]
        },
        {
            'name': '模組 8：整合與配置',
            'total': 10, 'p0': 5, 'p1': 3, 'p2': 2,
            'description': '與 EDR 系統整合，提供配置管理',
            'key_components': [
                'EDR 事件接收器（核心驅動、ETW、Hook）',
                '主工作流程（事件 → 評估 → 決策 → 響應）',
                '配置檔案解析（YAML 格式）',
                '配置管理（規則、效能、阻擋、告警、日誌）',
                '配置熱更新',
                '日誌記錄器（輪轉、過濾）'
            ],
            'challenges': [
                'EDR 介面對接',
                '配置驗證',
                '熱更新的一致性'
            ]
        },
        {
            'name': '模組 9：匹配引擎 API',
            'total': 7, 'p0': 4, 'p1': 3, 'p2': 0,
            'description': '提供完整的引擎 API 介面',
            'key_components': [
                'CapaEngine 類別（引擎主介面）',
                '生命週期管理（初始化、啟動、停止）',
                '事件處理 API（同步/非同步）',
                '規則管理 API（載入、卸載、啟用/停用）',
                '配置管理 API（執行時修改）',
                '統計與監控 API（指標、Prometheus）',
                '告警訂閱 API（回調註冊）',
                '白名單管理 API'
            ],
            'challenges': [
                'API 設計的易用性',
                '執行緒安全',
                '向後相容性'
            ]
        }
    ]

    for module in modules_detail:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = module['name']
        title.text_frame.paragraphs[0].font.size = Pt(28)

        # 模組資訊
        info_text = f"任務數：{module['total']} (P0:{module['p0']} | P1:{module['p1']} | P2:{module['p2']})\n{module['description']}"
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
        info_frame = info_box.text_frame
        info_frame.text = info_text
        info_frame.word_wrap = True
        for paragraph in info_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.name = '微軟正黑體'
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 51, 102)

        # 關鍵元件
        comp_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2.5))
        comp_frame = comp_box.text_frame
        comp_frame.word_wrap = True
        comp_frame.text = "關鍵元件："
        comp_frame.paragraphs[0].font.size = Pt(14)
        comp_frame.paragraphs[0].font.bold = True
        comp_frame.paragraphs[0].font.name = '微軟正黑體'

        for component in module['key_components']:
            p = comp_frame.add_paragraph()
            p.text = f"• {component}"
            p.font.size = Pt(12)
            p.font.name = '微軟正黑體'
            p.space_after = Pt(4)

        # 技術挑戰
        chal_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.8))
        chal_frame = chal_box.text_frame
        chal_frame.word_wrap = True
        chal_frame.text = "技術挑戰："
        chal_frame.paragraphs[0].font.size = Pt(14)
        chal_frame.paragraphs[0].font.bold = True
        chal_frame.paragraphs[0].font.name = '微軟正黑體'
        chal_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)

        for challenge in module['challenges']:
            p = chal_frame.add_paragraph()
            p.text = f"⚠ {challenge}"
            p.font.size = Pt(12)
            p.font.name = '微軟正黑體'
            p.space_after = Pt(4)

    # ========== 投影片 20: 效能需求 ==========
    headers = ['指標', '目標值', '最大可接受值']
    data = [
        ['單一事件處理延遲 (P95)', '< 10ms', '< 50ms'],
        ['規則評估延遲 (P95)', '< 5ms', '< 20ms'],
        ['告警生成延遲', '< 1ms', '< 5ms'],
        ['吞吐量（單核心）', '≥ 5,000 events/sec', '-'],
        ['吞吐量（4核心）', '≥ 15,000 events/sec', '-'],
        ['CPU 使用率（高負載）', '< 30%', '< 50%'],
        ['記憶體使用（高負載）', '< 1GB', '< 2GB']
    ]
    add_table_slide(prs, "效能需求與指標", headers, data)

    # ========== 投影片 21: 風險評估 ==========
    headers = ['風險類別', '風險描述', '影響', '對策']
    data = [
        ['技術風險', 'Hook 穩定性問題', '高', '多種 Hook 方案備援'],
        ['技術風險', '效能無法達標', '高', '持續效能優化、分階段實施'],
        ['技術風險', 'Windows 版本相容性', '中', '多版本測試環境'],
        ['業務風險', '誤報率過高', '高', '白名單機制、學習模式'],
        ['業務風險', '阻擋影響正常業務', '高', '多種阻擋模式、緊急關閉開關'],
        ['資源風險', '開發人力不足', '中', '外部資源補充、優先級調整'],
        ['時程風險', '需求變更', '中', '敏捷開發、分階段交付']
    ]
    add_table_slide(prs, "風險評估與對策", headers, data)

    # ========== 投影片 22: 資源需求 ==========
    add_content_slide(prs, "資源需求", {
        "人力資源": [
            "核心開發工程師：2-3 人（C++、Windows 系統開發）",
            "安全研究員：1 人（威脅分析、規則優化）",
            "測試工程師：1 人（效能測試、相容性測試）",
            "專案經理：1 人（兼職）"
        ],
        "開發環境": [
            "開發工具：Visual Studio 2022、Windows SDK",
            "建置系統：CMake、vcpkg",
            "測試環境：多版本 Windows（10/11/Server 2019/2022）",
            "CI/CD：GitHub Actions 或 Jenkins"
        ],
        "第三方函式庫": [
            "YAML 解析：yaml-cpp 或 RapidYAML",
            "JSON 解析：nlohmann/json 或 RapidJSON",
            "HTTP 通訊：libcurl 或 cpp-httplib",
            "加密：OpenSSL 或 Crypto++"
        ],
        "硬體資源": [
            "開發機：8 核心 CPU、32GB RAM、SSD",
            "測試機：多台虛擬機（不同 Windows 版本）",
            "效能測試機：16 核心 CPU、64GB RAM"
        ]
    })

    # ========== 投影片 23: 進度追蹤 ==========
    headers = ['模組', '任務數', 'P0', 'P1', 'P2', '完成數', '完成率']
    data = [
        ['模組 1：規則解析器', 12, 9, 2, 1, 0, '0%'],
        ['模組 2：事件抽象層', 12, 7, 4, 1, 0, '0%'],
        ['模組 3：匹配引擎', 31, 15, 13, 3, 0, '0%'],
        ['模組 4：決策與阻擋引擎', 8, 4, 3, 1, 0, '0%'],
        ['模組 5：響應動作執行器', 15, 9, 4, 2, 0, '0%'],
        ['模組 6：告警生成與管理', 9, 2, 6, 1, 0, '0%'],
        ['模組 7：效能監控', 9, 0, 7, 2, 0, '0%'],
        ['模組 8：整合與配置', 10, 5, 3, 2, 0, '0%'],
        ['模組 9：引擎 API', 7, 4, 3, 0, 0, '0%'],
        ['總計', 113, 55, 45, 13, 0, '0%']
    ]
    add_table_slide(prs, "進度追蹤（可定期更新）", headers, data)

    # ========== 投影片 24: Q&A ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    qa_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Q & A\n\n問題與討論"
    qa_para = qa_frame.paragraphs[0]
    qa_para.alignment = PP_ALIGN.CENTER
    qa_para.font.size = Pt(48)
    qa_para.font.bold = True
    qa_para.font.color.rgb = RGBColor(0, 51, 102)

    # 儲存簡報
    filename = f"EDR-Capa-詳細規格-{datetime.now().strftime('%Y%m%d')}.pptx"
    prs.save(filename)
    print(f"[OK] 簡報已成功建立：{filename}")
    print(f"總投影片數：{len(prs.slides)}")

if __name__ == "__main__":
    create_detailed_presentation()
